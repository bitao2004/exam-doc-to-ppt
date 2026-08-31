# -*- coding: utf-8 -*-
"""考题 docx → PPT（点击逐题揭示，WPS/PowerPoint 100% 兼容，默认精致美化）。

依赖（隔离 venv 安装，勿污染全局）：
    python -m venv <venv> && <venv>/Scripts/python.exe -m pip install python-pptx python-docx lxml

用法：
    python generate_ppt.py --src 考题.docx --out 考题.pptx --per-page 2 --reveal per-question

关键设计（踩坑沉淀）：
  * 绝不用手写动画时间线（WPS 会静默丢弃 → 答案直接显示）。
  * 用「幻灯片切换」实现点击显隐：每组生成多张物理幻灯片表达逐步揭示状态。
    - reveal=group      : 2 张/组（全隐 → 全显）
    - reveal=per-question: (per_page+1) 张/组（第0张全隐，每多一张多显1题，末张全显）
  * 美化（默认开启，--no-beautify 关闭）：渐变页眉+渐变背景、白色圆角卡片(柔影+左侧彩脊)、
    渐变圆形题号徽标、选项行圆角胶囊(字母蓝色圆标)、正确答案绿色高亮+✓、页脚提示药丸。
    布局按比例计算，答案/提示始终收在卡片内、不溢出。
"""
import argparse
import math
import re
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

EMU = 914400
FONT = "Microsoft YaHei"
def emu(inch): return int(inch * EMU)

# ---------- 版式（英寸，运行时转 EMU） ----------
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MX_IN = 0.55
HEADER_H_IN = 0.92
TOP_PAD_IN = 0.16
BOT_PAD_IN = 0.18
GAP_IN = 0.16
TITLE = "单选题"  # 默认页眉标题；main 中按文档自动识别或 --title 覆盖


def detect_title(src):
    """推断页眉标题。优先级：--title > 文件名关键字 > 文档正文关键字 > 默认"单选题"。

    题目文档首段通常是"1. 题干…"（形如题号），不能靠首段判断题型；
    改用文件名 / 正文中的"多选/判断/单选"关键字识别，最稳。
    """
    name = (src or "").lower()
    if "多选" in name:
        return "多选题"
    if "判断" in name:
        return "判断题"
    if "单选" in name:
        return "单选题"
    try:
        d = Document(src)
        full = "\n".join(p.text for p in d.paragraphs)
        if "多选" in full:
            return "多选题"
        if "判断" in full:
            return "判断题"
        if "单选" in full:
            return "单选题"
    except Exception:
        pass
    return "单选题"

# ---------- 配色（hex 字符串，便于改主题） ----------
# 背景：柔和浅色渐变（非深色纯色，避免刺眼）
BG            = "F7FAFE"   # 背景渐变 起（极浅冷白）
BG2           = "ECF1F9"   # 背景渐变 止（浅雾蓝）
# 页眉：柔和蓝紫渐变（比之前的深靛蓝更轻盈，白字可读且不刺眼）
HDR_C1        = "5B8FE6"   # 页眉渐变 起（柔和蓝）
HDR_C2        = "8E7BE8"   # 页眉渐变 止（柔和薰衣草）
CARD          = "FFFFFF"
CARD_BORDER   = "E6ECF3"
SPINE         = "5B8FE6"   # 卡片左侧彩脊
BADGE_C1      = "5B8FE6"   # 题号徽标渐变 起
BADGE_C2      = "8E7BE8"   # 题号徽标渐变 止
Q_TEXT        = "1F2937"
OPT_TEXT      = "3A4757"
OPT_ROW_BG    = "F7F9FC"
OPT_ROW_BORDER= "E6ECF3"
LETTER_BG     = "EDF2FD"
LETTER_FG     = "4F7BD6"
GREEN_BG      = "EAF8F0"
GREEN_BORDER  = "2DAE6B"
GREEN_FG      = "1F9D57"
CHECK         = "1F9D57"
HINT_PILL_BG  = "F0F3F8"
HINT_FG       = "6B7686"
ANS_HIDDEN_FG = "9AA6B6"
DIVIDER       = "E9EEF5"
SUBTITLE_FG   = "E8EEFF"
WHITE         = "FFFFFF"

# 页眉是否显示副标题文字（如"点击鼠标显示答案"）。
# 用户明确要求移除该类提示文字，故默认 False。
# 若将来希望恢复历史副标题提示，改 True 即可，build_header 会渲染 subtitle。
SHOW_SUBTITLE = False


# ---------- 1. 解析 ----------
def parse_questions(path):
    doc = Document(path)
    nonempty = [t.strip() for t in (p.text for p in doc.paragraphs) if t.strip()]
    ans_pat = re.compile(r'^\s*参考*答案\s*[:：]\s*([A-F]+)', re.UNICODE)
    q_pat = re.compile(r'^\s*(\d+)\s*[\.、]\s*(.*)')
    opt_pat = re.compile(r'^\s*([A-Fa-f])[\.、]\s*(.*)')
    questions, cur = [], None

    def flush():
        nonlocal cur
        if cur and cur.get('question'):
            questions.append(cur)
        cur = None

    for t in nonempty:
        if q_pat.match(t):
            flush(); cur = {'num': int(q_pat.match(t).group(1)),
                            'question': q_pat.match(t).group(2).strip(),
                            'options': [], 'answer': None}; continue
        mo = opt_pat.match(t)
        if mo and cur and cur['answer'] is None:
            cur['options'].append((mo.group(1).upper(), mo.group(2).strip())); continue
        ma = ans_pat.match(t)
        if ma and cur:
            cur['answer'] = ma.group(1).upper().replace(' ', ''); continue
        if cur:
            if cur['options']:
                cur['options'][-1] = (cur['options'][-1][0], cur['options'][-1][1] + ' ' + t)
            else:
                cur['question'] += ' ' + t
    flush()
    return questions


# ---------- 2. 样式 / 绘制 ----------
def set_cjk(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', name)

def style_run(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    set_cjk(run)

def add_shadow(shape, blur=72000, dist=36000, alpha=16):
    spPr = shape._element.spPr
    for e in spPr.findall(qn('a:effectLst')):
        spPr.remove(e)
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', str(blur)); sh.set('dist', str(dist))
    sh.set('dir', '5400000'); sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr')); clr.set('val', '1F2A44')
    a = etree.SubElement(clr, qn('a:alpha')); a.set('val', str(alpha * 1000))

def set_gradient(shape, c1, c2):
    try:
        shape.fill.gradient()
        stops = shape.fill.gradient_stops
        stops[0].position = 0.0
        stops[0].color.rgb = RGBColor.from_string(c1)
        stops[1].position = 1.0
        stops[1].color.rgb = RGBColor.from_string(c2)
    except Exception:
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor.from_string(c1)

def add_rect(slide, x, y, w, h, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             fill=None, line=None, line_w=1, radius=None, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        sp.fill.background()
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius is not None:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if shadow:
        add_shadow(sp)
    return sp

def add_text(slide, x, y, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    style_run(r, size, color, bold)
    return tb


# ---------- 3b. 字号自适应（保证长题干/长选项完整显示、不溢出卡片） ----------
def _char_units(ch):
    """估算单个字符的"宽度单位"（以 1pt 字号的单个全角字符 = 1.0 为单位）。"""
    o = ord(ch)
    if ch in ' \t\u3000':
        return 0.30
    if 0x2E80 <= o <= 0x9FFF or 0xF900 <= o <= 0xFAFF or 0xFF00 <= o <= 0xFFEF:
        return 1.0                      # CJK / 全角
    if ch.isascii():
        return 0.30 if (not ch.isalnum()) else 0.55  # 标点窄、字母数字中等
    return 0.6

def text_units(text):
    return sum(_char_units(c) for c in text)

def fit_font(text, box_w_in, box_h_in, max_f, min_f=9.0, line_ratio=1.2):
    """返回能完整装进 box（英寸）的最大字号（pt）。超长文本自动缩小到 min_f。"""
    w_pts = box_w_in * 72.0
    h_pts = box_h_in * 72.0
    units = text_units(text)
    if units <= 0:
        return float(max_f)
    f = float(max_f)
    while f >= min_f - 1e-9:
        cpl = max(1.0, w_pts / f)            # 每行可容纳的宽度单位
        lines = max(1, math.ceil(units / cpl))
        need = lines * f * line_ratio        # 所需高度（含行距）
        if need <= h_pts + 0.6:              # 容差，避免临界抖动
            return round(f, 1)
        f -= 0.5
    return round(min_f, 1)


# ---------- 3. 题块（精致美化，按卡片宽高比例布局，不溢出） ----------
def build_card(slide, q, left_in, top_in, cw_in, ch_in, per_page, show_ans, beautify):
    L = emu(left_in); T = emu(top_in)
    CW = emu(cw_in); CH = emu(ch_in)
    # 卡片底
    card = add_rect(slide, L, T, CW, CH, fill=CARD, line=CARD_BORDER, line_w=1,
                    radius=0.02, shadow=beautify)
    # 左侧彩脊
    if beautify:
        add_rect(slide, L, T + emu(0.12), emu(0.10), CH - emu(0.24),
                 fill=SPINE, radius=0.5)

    big = ch_in >= 2.6  # 卡片较高时用大字号，更舒展
    pad = 0.30
    # 题号徽标（圆形，直径取卡片宽高较小边的比例）
    badge = min(0.55, cw_in * 0.11, ch_in * 0.09)
    bx = L + emu(0.34); by = T + emu(0.30)
    bsp = add_rect(slide, bx, by, emu(badge), emu(badge), shape=MSO_SHAPE.OVAL,
                   fill=BADGE_C1, radius=None)
    if beautify:
        set_gradient(bsp, BADGE_C1, BADGE_C2)
    add_text(slide, bx, by, emu(badge), emu(badge), str(q['num']),
             18 if big else 14, WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 题干（徽标右侧，顶部与题号对齐；过长自动缩小字号，保证完整显示）
    qx = bx + emu(badge) + emu(0.16)
    qy = by  # 与题号徽标顶部对齐
    qw = CW - emu(badge) - emu(0.34 + 0.16) - emu(pad)
    qh = emu(ch_in * 0.24)
    qw_in = qw / EMU; qh_in = ch_in * 0.24
    qf = fit_font(q['question'], qw_in, qh_in, 18 if big else 15, min_f=10.5)
    add_text(slide, qx, qy, qw, qh, q['question'], qf, Q_TEXT,
             bold=True, anchor=MSO_ANCHOR.TOP)
    # 分隔线
    add_rect(slide, emu(left_in + 0.34), T + emu(ch_in * 0.30),
             emu(cw_in - 0.68), emu(0.012), fill=DIVIDER)

    # 选项（行距收紧，减小纵向占用；超长选项文字自动缩小字号）
    oy0 = T + emu(ch_in * 0.355)
    fy = T + CH - emu(0.56)       # 答案行 top
    ans_gap = emu(0.24)           # 答案行与最后一选项之间留间距
    opt_avail = (fy - ans_gap) - oy0
    n = max(1, len(q['options']))
    row_h = opt_avail * 0.90 / n
    g = opt_avail * 0.10 / (n - 1) if n > 1 else 0
    row_x = L + emu(0.30); row_w = CW - emu(0.60)
    row_w_in = row_w / EMU; row_h_in = row_h / EMU
    opt_maxf = 14 if big else 12
    for i, (letter, text) in enumerate(q['options']):
        ry = oy0 + i * (row_h + g)
        # 多选答案可为多字母（如 BCE）：用 in 判断以高亮所有正确选项
        is_ans = show_ans and letter in q['answer']
        rowfill = GREEN_BG if is_ans else OPT_ROW_BG
        rowline = GREEN_BORDER if is_ans else OPT_ROW_BORDER
        add_rect(slide, row_x, ry, row_w, row_h, fill=rowfill, line=rowline,
                 line_w=1, radius=0.18)
        lb = min(emu(0.34), emu((row_h / EMU) * 0.6))
        lx = row_x + emu(0.16); ly = ry + (row_h - lb) // 2
        lfill = GREEN_FG if is_ans else LETTER_BG
        lfg = WHITE if is_ans else LETTER_FG
        add_rect(slide, lx, ly, lb, lb, shape=MSO_SHAPE.OVAL, fill=lfill)
        add_text(slide, lx, ly, lb, lb, letter, 14 if big else 11, lfg, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        of = fit_font(text, row_w_in - 1.10, row_h_in, opt_maxf, min_f=9.0)
        add_text(slide, row_x + emu(0.70), ry, row_w - emu(1.20), row_h, text,
                 of, OPT_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        if is_ans:
            add_text(slide, row_x + row_w - emu(0.78), ry, emu(0.66), row_h, '✓',
                     18 if big else 15, CHECK, bold=True, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    # 页脚：提示药丸 + 答案
    pill_w = emu(2.6)
    add_rect(slide, row_x, fy, pill_w, emu(0.36), fill=HINT_PILL_BG, radius=0.5)
    if show_ans:
        pill_text = "✅ 答案已显示"; pill_color = GREEN_FG
        ans_text = f"参考答案：{q['answer']}"; ans_color = GREEN_FG; ans_bold = True
    else:
        pill_text = "👉 点击显示答案"; pill_color = HINT_FG
        ans_text = "参考答案：？"; ans_color = ANS_HIDDEN_FG; ans_bold = False
    add_text(slide, row_x + emu(0.1), fy, pill_w - emu(0.2), emu(0.36),
             pill_text, 12 if big else 10, pill_color, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, row_x + row_w - emu(3.0), fy, emu(2.8), emu(0.36),
             ans_text, 13 if big else 11, ans_color, bold=ans_bold,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def build_header(slide, group_idx, total_groups, subtitle, beautify):
    # 背景：柔和浅色渐变（非深色纯色，避免刺眼）
    bg = add_rect(slide, 0, 0, emu(SLIDE_W_IN), emu(SLIDE_H_IN),
                  shape=MSO_SHAPE.RECTANGLE, fill=BG)
    if beautify:
        set_gradient(bg, BG, BG2)
    # 页眉条（柔和蓝紫渐变）
    hb = add_rect(slide, 0, 0, emu(SLIDE_W_IN), emu(HEADER_H_IN),
                  shape=MSO_SHAPE.RECTANGLE, fill=HDR_C2)
    if beautify:
        set_gradient(hb, HDR_C1, HDR_C2)
    add_text(slide, emu(MX_IN), emu(0.14), emu(6), emu(0.42), TITLE, 28, WHITE,
             bold=True, anchor=MSO_ANCHOR.MIDDLE)
    pw = emu(2.3); px = emu(SLIDE_W_IN - MX_IN) - pw; py = emu(0.26); ph = emu(0.40)
    add_rect(slide, px, py, pw, ph, fill=WHITE, radius=0.5)
    add_text(slide, px, py, pw, ph, f"第 {group_idx} / {total_groups} 组", 13, HDR_C1,
             bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 副标题文字（默认不渲染，见 SHOW_SUBTITLE；恢复时显示在标题与页码之间）
    if SHOW_SUBTITLE and subtitle:
        add_text(slide, emu(6.6), emu(0.14), emu(3.6), emu(0.42), subtitle, 12,
                 SUBTITLE_FG, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ---------- 4. 揭示状态 ----------
def reveal_states(per_page, mode):
    if mode == "group":
        return [[False] * per_page, [True] * per_page]
    return [[j < k for j in range(per_page)] for k in range(per_page + 1)]

def subtitle_for(k, per_page, mode):
    if mode == "group":
        return "点击鼠标显示答案" if k == 0 else "点击继续下一页"
    if k == 0:
        return "点击鼠标显示答案"
    if k < per_page:
        return "点击显示下一题答案"
    return "点击继续下一页"


def build_group(prs, blank, pair, group_no, total_groups, per_page, mode, beautify, layout):
    states = reveal_states(per_page, mode)
    # 左右双栏：每页 per_page 张卡并排，单卡接近方形，充分利用宽屏
    if layout == "horizontal" and per_page > 1:
        gap = GAP_IN
        cw_in = (SLIDE_W_IN - 2 * MX_IN - (per_page - 1) * gap) / per_page
        ch_in = SLIDE_H_IN - HEADER_H_IN - TOP_PAD_IN - BOT_PAD_IN
        top_in = HEADER_H_IN + TOP_PAD_IN
        for k, state in enumerate(states):
            s = prs.slides.add_slide(blank)
            build_header(s, group_no, total_groups, subtitle_for(k, per_page, mode), beautify)
            for i, q in enumerate(pair):
                left_in = MX_IN + i * (cw_in + gap)
                build_card(s, q, left_in, top_in, cw_in, ch_in, per_page, state[i], beautify)
    else:
        # 上下堆叠（默认）
        CH_in = (SLIDE_H_IN - HEADER_H_IN - TOP_PAD_IN - BOT_PAD_IN - (per_page - 1) * GAP_IN) / per_page
        CH = emu(CH_in)
        CT1 = emu(HEADER_H_IN + TOP_PAD_IN)
        for k, state in enumerate(states):
            s = prs.slides.add_slide(blank)
            build_header(s, group_no, total_groups, subtitle_for(k, per_page, mode), beautify)
            for i, q in enumerate(pair):
                top = CT1 + i * (CH + emu(GAP_IN))
                build_card(s, q, MX_IN, top / EMU, SLIDE_W_IN - 2 * MX_IN, CH_in,
                           per_page, state[i], beautify)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--src", default="考题.docx")
    ap.add_argument("--out", default="考题.pptx")
    ap.add_argument("--title", default=None,
                    help="页眉标题，默认按文档首段自动识别（单选题/多选题/判断题…）")
    ap.add_argument("--per-page", type=int, default=2)
    ap.add_argument("--reveal", choices=["group", "per-question"], default="per-question")
    ap.add_argument("--layout", choices=["vertical", "horizontal"], default="vertical",
                    help="vertical=上下堆叠; horizontal=左右双栏(充分利用宽屏)")
    ap.add_argument("--beautify", dest="beautify", action="store_true", default=True)
    ap.add_argument("--no-beautify", dest="beautify", action="store_false")
    args = ap.parse_args()
    global TITLE
    TITLE = args.title or (detect_title(args.src) or "单选题")

    questions = parse_questions(args.src)
    print("解析题目数:", len(questions))

    prs = Presentation()
    prs.slide_width = Emu(emu(SLIDE_W_IN))
    prs.slide_height = Emu(emu(SLIDE_H_IN))
    blank = prs.slide_layouts[6]
    total_groups = (len(questions) + args.per_page - 1) // args.per_page
    for i in range(0, len(questions), args.per_page):
        pair = questions[i:i + args.per_page]
        build_group(prs, blank, pair, i // args.per_page + 1, total_groups,
                    args.per_page, args.reveal, args.beautify, args.layout)

    candidates = [args.out, args.out[:-5] + "_v2.pptx", "考题_生成.pptx"]
    saved = None
    for c in candidates:
        try:
            prs.save(c); saved = c; break
        except OSError:
            continue
    if saved is None:
        saved = "考题_生成.pptx"; prs.save(saved)
    if saved != args.out:
        # 文件锁定坑：目标被预览面板/WPS 占用时写入会 Permission denied，
        # 已自动改存备用名——必须显式提醒用户，切勿静默丢结果。
        print(f"⚠️ 目标文件 {args.out} 被占用（可能正被预览/WPS 打开），已另存为：{saved}")
        print("   请关闭该文件的预览或程序后，手动将其重命名/覆盖为期望文件名。")
    print("已保存:", saved, " 物理幻灯片数:", len(prs.slides._sldIdLst),
          " 美化:", args.beautify)


if __name__ == "__main__":
    main()
